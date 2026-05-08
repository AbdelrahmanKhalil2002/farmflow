"""
FarmFlow PowerPoint Generator
Creates two presentations: English and Arabic
Run: python3 generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ─── Brand Colors ────────────────────────────────────────────────
GREEN_DARK   = RGBColor(0x1B, 0x5E, 0x20)   # #1B5E20
GREEN_MID    = RGBColor(0x2E, 0x7D, 0x32)   # #2E7D32
GREEN_LIGHT  = RGBColor(0x43, 0xA0, 0x47)   # #43A047
GREEN_PALE   = RGBColor(0xE8, 0xF5, 0xE9)   # #E8F5E9
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x21, 0x21, 0x21)    # #212121
GREY_TEXT    = RGBColor(0x61, 0x61, 0x61)    # #616161
ACCENT       = RGBColor(0xFF, 0xB3, 0x00)    # #FFB300 amber accent

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── Helpers ─────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def txb(slide, text, x, y, w, h,
        font_size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, rtl=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    if rtl:
        _set_rtl(p)
    return tb


def _set_rtl(paragraph):
    """Set RTL on a paragraph XML element."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set(qn('a:rtl'), '1')


def multi_para_txb(slide, lines, x, y, w, h,
                   font_size=16, bold_first=False,
                   text_color=DARK_TEXT, bullet=False,
                   align=PP_ALIGN.LEFT, rtl=False):
    """Add a textbox with multiple paragraphs. lines = list of strings."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if rtl:
            _set_rtl(p)
        if bullet and not first:
            p.level = 1
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and (line == lines[0])
        run.font.color.rgb = text_color
    return tb


def divider(slide, y, color=GREEN_MID):
    rect(slide, Inches(0.4), y, Inches(12.5), Pt(2), color)


# ─── ENGLISH SLIDES ──────────────────────────────────────────────

def slide_cover_en(prs):
    s = blank(prs)
    # Full background
    rect(s, 0, 0, SLIDE_W, SLIDE_H, GREEN_DARK)
    # Accent strip
    rect(s, 0, Inches(5.8), SLIDE_W, Inches(1.7), GREEN_MID)
    # Decorative circle (top right)
    c = s.shapes.add_shape(9, Inches(10.5), Inches(-1), Inches(4), Inches(4))  # oval
    c.fill.solid(); c.fill.fore_color.rgb = GREEN_MID
    c.line.fill.background()
    # Title
    txb(s, "🌿 FarmFlow", Inches(0.7), Inches(1.5), Inches(10), Inches(1.2),
        font_size=56, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Tagline
    txb(s, "Egypt's Livestock Marketplace — Digital, Trusted, Modern",
        Inches(0.7), Inches(2.9), Inches(11), Inches(0.7),
        font_size=22, color=RGBColor(0xC8, 0xE6, 0xC9), align=PP_ALIGN.LEFT)
    # Sub
    txb(s, "Connecting Farm Owners with Buyers across Egypt",
        Inches(0.7), Inches(3.6), Inches(11), Inches(0.6),
        font_size=17, color=RGBColor(0xA5, 0xD6, 0xA7), italic=True)
    # Bottom bar text
    txb(s, "Web  ·  Mobile (iOS & Android)  ·  Desktop",
        Inches(0.7), Inches(6.1), Inches(10), Inches(0.6),
        font_size=16, color=WHITE)
    txb(s, "Version 1.0  —  May 2026",
        Inches(0.7), Inches(6.7), Inches(10), Inches(0.5),
        font_size=12, color=RGBColor(0xA5, 0xD6, 0xA7), italic=True)


def slide_problem_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "The Problem", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        font_size=34, bold=True, color=WHITE)
    problems = [
        ("❌  No centralized marketplace",     "Buyers have nowhere to discover all farms at once"),
        ("❌  No animal health transparency",   "Buyers can't verify the health of an animal before purchase"),
        ("❌  No herd management tools",         "Farmers track animals in notebooks or memory"),
        ("❌  No financial tracking",            "Expenses and income are never properly recorded"),
        ("❌  Informal, unverified deals",       "No accountability or verification process exists"),
    ]
    y = Inches(1.5)
    for title, sub in problems:
        rect(s, Inches(0.4), y, Inches(12.4), Inches(0.85), GREEN_PALE)
        txb(s, title, Inches(0.6), y + Pt(4), Inches(7), Inches(0.4),
            font_size=15, bold=True, color=GREEN_DARK)
        txb(s, sub, Inches(7.2), y + Pt(4), Inches(5.5), Inches(0.4),
            font_size=13, color=GREY_TEXT)
        y += Inches(1.0)


def slide_solution_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_MID)
    txb(s, "The Solution — FarmFlow", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        font_size=34, bold=True, color=WHITE)
    txb(s, "A single digital platform that manages the entire livestock trade lifecycle",
        Inches(0.5), Inches(1.45), Inches(12), Inches(0.5),
        font_size=17, color=GREY_TEXT)
    # 3 columns
    cols = [
        (GREEN_DARK,  "🐄  Farm Owners",  ["Manage every animal", "Track vaccinations", "Record medical history", "Log expenses & income", "List animals for sale", "Export financial reports"]),
        (GREEN_MID,   "🛒  Buyers",        ["Browse all farms", "See health certificates", "View photos & videos", "Chat via WhatsApp", "Place orders online", "Rate & review sellers"]),
        (GREEN_LIGHT, "🛡  Platform",      ["Admin approval workflow", "Verified seller profiles", "Eid season mode", "Multi-product: livestock,\ndairy & supplies", "Push notifications", "Offline desktop app"]),
    ]
    for i, (color, title, bullets) in enumerate(cols):
        x = Inches(0.4 + i * 4.3)
        rect(s, x, Inches(2.1), Inches(4.0), Inches(4.9), color)
        txb(s, title, x + Inches(0.15), Inches(2.2), Inches(3.7), Inches(0.55),
            font_size=17, bold=True, color=WHITE)
        y = Inches(2.85)
        for b in bullets:
            txb(s, f"• {b}", x + Inches(0.15), y, Inches(3.7), Inches(0.5),
                font_size=13, color=WHITE)
            y += Inches(0.55)


def slide_platforms_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "Available On Every Platform", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        font_size=34, bold=True, color=WHITE)
    platforms = [
        ("🌐", "Web App",        "Browser",          "Online",               "Sellers & Buyers",      GREEN_DARK),
        ("📱", "Mobile App",     "iOS & Android",    "Online",               "Sellers & Buyers",      GREEN_MID),
        ("🖥️", "Desktop App",    "Windows & macOS",  "Offline + Online ✓",   "Farm Owners (PCs)",     GREEN_LIGHT),
    ]
    y = Inches(1.55)
    for icon, name, os_, conn, users, color in platforms:
        rect(s, Inches(0.4), y, Inches(12.5), Inches(1.5), color)
        txb(s, icon + "  " + name, Inches(0.6), y + Inches(0.25), Inches(3), Inches(0.9),
            font_size=22, bold=True, color=WHITE)
        txb(s, os_, Inches(3.7), y + Inches(0.1), Inches(2.5), Inches(0.55),
            font_size=15, color=WHITE)
        txb(s, conn, Inches(6.3), y + Inches(0.1), Inches(2.8), Inches(0.55),
            font_size=15, bold=True, color=ACCENT)
        txb(s, users, Inches(9.2), y + Inches(0.1), Inches(3.5), Inches(0.55),
            font_size=15, color=WHITE, italic=True)
        y += Inches(1.65)
    txb(s, "★  The Desktop App works fully OFFLINE — perfect for farm owners in areas with poor connectivity",
        Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5),
        font_size=13, color=GREEN_DARK, italic=True)


def slide_seller_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "Seller Features — Farm Owner", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=34, bold=True, color=WHITE)
    features = [
        ("📊  Dashboard",         "KPIs, income chart, vaccination & weighing reminders, quick actions"),
        ("🐄  Herd Management",   "Register every animal: breed, age, weight, health, pregnancy, photos"),
        ("💉  Vaccinations",      "Log vaccination records; color-coded overdue/upcoming alerts"),
        ("🏥  Medical Records",   "Vet visits with diagnosis, medication, cost — auto-linked to expenses"),
        ("📋  Listings",          "Only animals you choose appear to buyers; full control"),
        ("💰  Finance",           "Expense categories, income tracking, bar charts, CSV/PDF export"),
        ("🥛  Dairy Products",    "Milk, cheese, butter, ghee — list & manage separately"),
        ("🌾  Farm Supplies",     "Feed, equipment, seeds — additional revenue stream"),
        ("👤  Seller Profile",    "Farm banner, bio, certifications, contact details"),
    ]
    y = Inches(1.4)
    col_break = 5
    for i, (feat, desc) in enumerate(features):
        if i < col_break:
            x = Inches(0.4)
        else:
            x = Inches(6.8)
            if i == col_break:
                y = Inches(1.4)
        row_y = y if i >= col_break else Inches(1.4) + Inches((i) * 1.1)
        if i >= col_break:
            row_y = Inches(1.4) + Inches((i - col_break) * 1.1)
        rect(s, x, row_y, Inches(6.1), Inches(0.95), GREEN_PALE)
        txb(s, feat, x + Inches(0.1), row_y + Pt(4), Inches(2.2), Inches(0.45),
            font_size=13, bold=True, color=GREEN_DARK)
        txb(s, desc, x + Inches(2.3), row_y + Pt(4), Inches(3.8), Inches(0.45),
            font_size=12, color=DARK_TEXT)


def slide_herd_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_MID)
    txb(s, "Herd Management — Complete Animal Lifecycle", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=30, bold=True, color=WHITE)
    # Animal profile tabs
    tabs = [
        (GREEN_DARK,  "📈 Growth",       ["Weight history chart", "Set target weight", "Track progress over time", "Log new weight entries"]),
        (GREEN_MID,   "💉 Vaccinations",  ["Log all vaccinations", "Due date alerts", "Overdue color-coding", "Upcoming vaccine reminders"]),
        (GREEN_LIGHT, "🏥 Medical",       ["Vet visit records", "Diagnosis & treatment", "Medication details", "Auto-links cost to expenses"]),
        (RGBColor(0x1B,0x5E,0x20), "ℹ️ Info", ["Breed & tag ID", "Date of birth", "Gender & color", "Pregnancy status + due date"]),
    ]
    for i, (color, title, items) in enumerate(tabs):
        x = Inches(0.4 + i * 3.2)
        rect(s, x, Inches(1.5), Inches(3.0), Inches(4.9), color)
        txb(s, title, x + Inches(0.1), Inches(1.6), Inches(2.8), Inches(0.55),
            font_size=17, bold=True, color=WHITE)
        y = Inches(2.25)
        for item in items:
            txb(s, f"• {item}", x + Inches(0.1), y, Inches(2.8), Inches(0.5),
                font_size=13, color=WHITE)
            y += Inches(0.55)
    txb(s, "🔒  Non-sale animals stay PRIVATE — only animals you choose are visible to buyers",
        Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.55),
        font_size=14, bold=True, color=GREEN_DARK)


def slide_finance_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "Financial Management — Built for Farm Owners", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=30, bold=True, color=WHITE)
    # Left: expense categories
    rect(s, Inches(0.4), Inches(1.4), Inches(4.1), Inches(5.7), GREEN_PALE)
    txb(s, "Expense Categories", Inches(0.5), Inches(1.5), Inches(3.9), Inches(0.5),
        font_size=15, bold=True, color=GREEN_DARK)
    cats = ["🌾 Feed & Fodder", "💊 Veterinary", "💧 Water", "⚡ Electricity",
            "👷 Employee Salaries", "🏠 Farm Rent", "🚛 Transport", "🔧 Maintenance", "📦 Other"]
    y = Inches(2.1)
    for c in cats:
        txb(s, c, Inches(0.6), y, Inches(3.7), Inches(0.42),
            font_size=13, color=DARK_TEXT)
        y += Inches(0.5)
    # Right: reports
    reports = [
        ("📊  Bar Chart Reports",   "Income vs. expenses by month or year"),
        ("🔍  Drill-Down",          "Tap any bar to see individual transactions"),
        ("📋  Full History",        "Categorized expense & income lists"),
        ("📁  CSV Export",          "Share data as a spreadsheet"),
        ("📄  PDF Export",          "Print-ready A4 financial report"),
        ("🎯  Budget Tracking",     "Set monthly/yearly budget per category"),
        ("💹  KPI Strip",           "Total income, expenses, net profit, best month"),
    ]
    y = Inches(1.4)
    for title, desc in reports:
        rect(s, Inches(4.75), y, Inches(8.15), Inches(0.72), GREEN_PALE)
        txb(s, title, Inches(4.9), y + Pt(4), Inches(2.5), Inches(0.38),
            font_size=13, bold=True, color=GREEN_DARK)
        txb(s, desc, Inches(7.5), y + Pt(4), Inches(5.3), Inches(0.38),
            font_size=12, color=DARK_TEXT)
        y += Inches(0.82)


def slide_buyer_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_MID)
    txb(s, "Buyer Features — Find & Order Livestock", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=34, bold=True, color=WHITE)
    features = [
        ("🔍  Discover",       "Browse all farms; filter by governorate, animal type, price, ratings"),
        ("🏡  Farm Profile",   "See banner, ratings, contact info, all available livestock & dairy"),
        ("📸  Rich Listings",  "Photos, video, weight, age, breed, health certificate, price"),
        ("❤️  Health Cert",    "Every animal shows verified health status — know what you're buying"),
        ("📞  Contact Seller", "WhatsApp chat, direct call, or copy phone number — one tap"),
        ("🛒  Place Order",    "Delivery or pickup; Cash on Delivery or InstaPay"),
        ("📦  Track Orders",   "Real-time order status: pending → confirmed → in transit → delivered"),
        ("⭐  Reviews",        "Rate and review sellers after delivery; read others' reviews first"),
        ("🔔  Notifications",  "Instant alerts for order updates and special offers"),
        ("❤️  Favorites",      "Save farms to a favorites list for quick access"),
    ]
    y = Inches(1.4)
    col_break = 5
    for i, (feat, desc) in enumerate(features):
        col = 0 if i < col_break else 1
        xi = Inches(0.4) if col == 0 else Inches(6.8)
        row_y = Inches(1.4) + Inches((i % col_break) * 1.1)
        rect(s, xi, row_y, Inches(6.1), Inches(0.95), GREEN_PALE)
        txb(s, feat, xi + Inches(0.1), row_y + Pt(4), Inches(2.1), Inches(0.45),
            font_size=13, bold=True, color=GREEN_DARK)
        txb(s, desc, xi + Inches(2.2), row_y + Pt(4), Inches(3.8), Inches(0.45),
            font_size=12, color=DARK_TEXT)


def slide_coming_soon_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "Coming Soon", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        font_size=34, bold=True, color=WHITE)
    items = [
        ("🔪  Butcher On-Demand",          "Request a certified butcher to come to your location — no travel needed"),
        ("🥩  Ready-Slaughtered Delivery",  "Receive a slaughtered & prepared animal delivered to your door"),
        ("🔨  Live Livestock Auctions",     "Bid on premium animals in real-time online auctions"),
        ("🤖  AI Pricing Assistant",        "Smart price recommendations based on weight, breed & market trends"),
        ("🌿  Multi-Branch Management",     "Manage multiple farm locations from one seller account"),
        ("💳  Card Payment Gateway",        "Full Fawry / PayMob integration for Visa & Mastercard"),
    ]
    y = Inches(1.45)
    for icon_title, desc in items:
        rect(s, Inches(0.4), y, Inches(12.5), Inches(0.82), GREEN_PALE)
        txb(s, icon_title, Inches(0.6), y + Pt(5), Inches(3.5), Inches(0.4),
            font_size=15, bold=True, color=GREEN_DARK)
        txb(s, desc, Inches(4.2), y + Pt(5), Inches(8.6), Inches(0.4),
            font_size=13, color=DARK_TEXT)
        y += Inches(0.92)


def slide_closing_en(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, GREEN_DARK)
    rect(s, 0, Inches(5.5), SLIDE_W, Inches(2.0), GREEN_MID)
    # Decorative
    c = s.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(5), Inches(5))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN_MID; c.line.fill.background()
    txb(s, "Thank You", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
        font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, "FarmFlow — Bringing Egypt's livestock market into the digital age",
        Inches(1), Inches(3.2), Inches(11), Inches(0.8),
        font_size=20, color=RGBColor(0xC8, 0xE6, 0xC9), align=PP_ALIGN.CENTER, italic=True)
    txb(s, "Web  ·  Mobile  ·  Desktop  ·  Offline-capable",
        Inches(1), Inches(5.8), Inches(11), Inches(0.6),
        font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, "Version 1.0  —  May 2026",
        Inches(1), Inches(6.6), Inches(11), Inches(0.5),
        font_size=12, color=RGBColor(0xA5, 0xD6, 0xA7),
        align=PP_ALIGN.CENTER, italic=True)


# ─── ARABIC SLIDES ───────────────────────────────────────────────

def slide_cover_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, GREEN_DARK)
    rect(s, 0, Inches(5.8), SLIDE_W, Inches(1.7), GREEN_MID)
    c = s.shapes.add_shape(9, Inches(10.5), Inches(-1), Inches(4), Inches(4))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN_MID; c.line.fill.background()
    txb(s, "فارم فلو 🌿", Inches(0.7), Inches(1.5), Inches(11), Inches(1.2),
        font_size=56, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    txb(s, "سوق المواشي المصري — رقمي، موثوق، عصري",
        Inches(0.7), Inches(2.9), Inches(11), Inches(0.7),
        font_size=22, color=RGBColor(0xC8, 0xE6, 0xC9), align=PP_ALIGN.RIGHT, rtl=True)
    txb(s, "يربط أصحاب المزارع بالمشترين في جميع أنحاء مصر",
        Inches(0.7), Inches(3.6), Inches(11), Inches(0.6),
        font_size=17, color=RGBColor(0xA5, 0xD6, 0xA7), italic=True,
        align=PP_ALIGN.RIGHT, rtl=True)
    txb(s, "ويب  ·  جوال (iOS & Android)  ·  سطح المكتب",
        Inches(0.7), Inches(6.1), Inches(11.5), Inches(0.6),
        font_size=16, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    txb(s, "الإصدار 1.0  —  مايو 2026",
        Inches(0.7), Inches(6.7), Inches(11.5), Inches(0.5),
        font_size=12, color=RGBColor(0xA5, 0xD6, 0xA7), italic=True,
        align=PP_ALIGN.RIGHT, rtl=True)


def slide_problem_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "المشكلة", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=34, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    problems = [
        ("❌  لا يوجد سوق مركزي موحد",          "المشترون لا يجدون مكاناً واحداً يكتشفون فيه جميع المزارع"),
        ("❌  غياب الشفافية عن صحة الحيوانات",  "المشترون لا يستطيعون التحقق من صحة الحيوان قبل الشراء"),
        ("❌  لا توجد أدوات لإدارة القطيع",      "المزارعون يتتبعون الحيوانات في دفاتر أو من الذاكرة"),
        ("❌  لا يوجد تتبع مالي",                "المصاريف والدخل لا تُسجَّل بشكل صحيح أبداً"),
        ("❌  صفقات غير رسمية وغير موثقة",       "لا يوجد نظام للمحاسبة أو التحقق"),
    ]
    y = Inches(1.5)
    for title, sub in problems:
        rect(s, Inches(0.4), y, Inches(12.4), Inches(0.85), GREEN_PALE)
        txb(s, title, Inches(0.6), y + Pt(4), Inches(5.5), Inches(0.4),
            font_size=15, bold=True, color=GREEN_DARK, align=PP_ALIGN.RIGHT, rtl=True)
        txb(s, sub, Inches(0.7), y + Pt(4), Inches(11.5), Inches(0.4),
            font_size=13, color=GREY_TEXT, align=PP_ALIGN.RIGHT, rtl=True)
        y += Inches(1.0)


def slide_solution_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_MID)
    txb(s, "الحل — فارم فلو", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=34, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    txb(s, "منصة رقمية واحدة تدير دورة حياة تجارة المواشي بالكامل",
        Inches(0.5), Inches(1.45), Inches(12), Inches(0.5),
        font_size=17, color=GREY_TEXT, align=PP_ALIGN.RIGHT, rtl=True)
    cols = [
        (GREEN_DARK,  "🐄  التجار",     ["إدارة كل حيوان", "تتبع التطعيمات", "السجل الطبي", "تسجيل المصاريف والدخل", "عرض الحيوانات للبيع", "تصدير تقارير مالية"]),
        (GREEN_MID,   "🛒  المشترون",   ["تصفح جميع المزارع", "مشاهدة الشهادات الصحية", "صور وفيديو", "التواصل عبر واتساب", "تقديم الطلبات أونلاين", "تقييم البائعين"]),
        (GREEN_LIGHT, "🛡  المنصة",     ["سير عمل موافقة إدارية", "ملفات تعريف موثقة", "وضع العيد", "مواشي وألبان ومستلزمات", "إشعارات فورية", "تطبيق أوفلاين"]),
    ]
    for i, (color, title, bullets) in enumerate(cols):
        x = Inches(0.4 + i * 4.3)
        rect(s, x, Inches(2.1), Inches(4.0), Inches(4.9), color)
        txb(s, title, x + Inches(0.15), Inches(2.2), Inches(3.7), Inches(0.55),
            font_size=17, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
        y = Inches(2.85)
        for b in bullets:
            txb(s, f"• {b}", x + Inches(0.15), y, Inches(3.7), Inches(0.5),
                font_size=13, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
            y += Inches(0.55)


def slide_platforms_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "متاح على كل المنصات", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=34, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    platforms = [
        ("🌐", "تطبيق الويب",        "المتصفح",          "أونلاين فقط",               "البائعون والمشترون",   GREEN_DARK),
        ("📱", "تطبيق الجوال",       "iOS و Android",    "أونلاين",                    "البائعون والمشترون",   GREEN_MID),
        ("🖥️", "تطبيق سطح المكتب",  "Windows و macOS",  "✓ أوفلاين + أونلاين",        "أصحاب المزارع",        GREEN_LIGHT),
    ]
    y = Inches(1.55)
    for icon, name, os_, conn, users, color in platforms:
        rect(s, Inches(0.4), y, Inches(12.5), Inches(1.5), color)
        txb(s, icon + "  " + name, Inches(0.6), y + Inches(0.25), Inches(4.5), Inches(0.9),
            font_size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
        txb(s, os_, Inches(0.6), y + Inches(0.1), Inches(12), Inches(0.55),
            font_size=15, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
        txb(s, conn, Inches(5), y + Inches(0.7), Inches(7.5), Inches(0.55),
            font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT, rtl=True)
        y += Inches(1.65)
    txb(s, "★  تطبيق سطح المكتب يعمل بكامل وظائفه بدون إنترنت — مثالي لأصحاب المزارع في المناطق بعيدة",
        Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5),
        font_size=13, color=GREEN_DARK, italic=True, align=PP_ALIGN.RIGHT, rtl=True)


def slide_seller_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "مميزات التاجر — صاحب المزرعة", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=34, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    features = [
        ("📊  لوحة التحكم",      "مؤشرات الأداء، مخطط الدخل، تذكيرات التطعيم والوزن"),
        ("🐄  إدارة القطيع",     "تسجيل كل حيوان: السلالة، العمر، الوزن، الصحة، الحمل"),
        ("💉  التطعيمات",        "سجلات التطعيم مع تنبيهات ملونة للمتأخر والقادم"),
        ("🏥  السجل الطبي",      "زيارات البيطري مع التشخيص والعلاج والتكلفة"),
        ("📋  الإعلانات",        "فقط الحيوانات التي تختارها تظهر للمشترين"),
        ("💰  الإدارة المالية",   "فئات المصاريف، تتبع الدخل، مخططات، تصدير CSV/PDF"),
        ("🥛  منتجات الألبان",   "حليب، جبن، زبدة، سمن — قسم منفصل للإدارة"),
        ("🌾  مستلزمات المزرعة", "أعلاف، معدات، بذور — مصدر دخل إضافي"),
        ("👤  ملف التاجر",       "غلاف المزرعة، السيرة الذاتية، بيانات الاتصال"),
    ]
    y = Inches(1.4)
    for i, (feat, desc) in enumerate(features):
        col = 0 if i < 5 else 1
        xi = Inches(0.4) if col == 0 else Inches(6.8)
        row_y = Inches(1.4) + Inches((i % 5) * 1.1)
        rect(s, xi, row_y, Inches(6.1), Inches(0.95), GREEN_PALE)
        txb(s, feat, xi + Inches(0.1), row_y + Pt(4), Inches(5.9), Inches(0.45),
            font_size=13, bold=True, color=GREEN_DARK, align=PP_ALIGN.RIGHT, rtl=True)
        txb(s, desc, xi + Inches(0.1), row_y + Pt(24), Inches(5.9), Inches(0.38),
            font_size=12, color=DARK_TEXT, align=PP_ALIGN.RIGHT, rtl=True)


def slide_herd_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_MID)
    txb(s, "إدارة القطيع — دورة حياة كاملة لكل حيوان", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    tabs = [
        (GREEN_DARK,  "📈 النمو",         ["مخطط سجل الوزن", "تحديد الوزن المستهدف", "تتبع التقدم", "تسجيل أوزان جديدة"]),
        (GREEN_MID,   "💉 التطعيمات",     ["تسجيل التطعيمات", "تنبيهات الاستحقاق", "ترميز لوني للمتأخر", "تذكيرات القادم"]),
        (GREEN_LIGHT, "🏥 الطبي",         ["سجلات البيطري", "التشخيص والعلاج", "الأدوية", "التكلفة مرتبطة تلقائياً"]),
        (RGBColor(0x1B,0x5E,0x20), "ℹ️ المعلومات", ["السلالة ورقم الوسم", "تاريخ الميلاد", "الجنس واللون", "حالة الحمل وتاريخ الوضع"]),
    ]
    for i, (color, title, items) in enumerate(tabs):
        x = Inches(0.4 + i * 3.2)
        rect(s, x, Inches(1.5), Inches(3.0), Inches(4.9), color)
        txb(s, title, x + Inches(0.1), Inches(1.6), Inches(2.8), Inches(0.55),
            font_size=17, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
        y = Inches(2.25)
        for item in items:
            txb(s, f"• {item}", x + Inches(0.1), y, Inches(2.8), Inches(0.5),
                font_size=13, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
            y += Inches(0.55)
    txb(s, "🔒  الحيوانات الغير معروضة تبقى خاصة — فقط ما تختاره يظهر للمشترين",
        Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.55),
        font_size=14, bold=True, color=GREEN_DARK, align=PP_ALIGN.RIGHT, rtl=True)


def slide_finance_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "الإدارة المالية — مصممة لأصحاب المزارع", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    rect(s, Inches(9.1), Inches(1.4), Inches(3.9), Inches(5.7), GREEN_PALE)
    txb(s, "فئات المصاريف", Inches(9.2), Inches(1.5), Inches(3.7), Inches(0.5),
        font_size=15, bold=True, color=GREEN_DARK, align=PP_ALIGN.RIGHT, rtl=True)
    cats = ["🌾 الأعلاف", "💊 البيطري", "💧 المياه", "⚡ الكهرباء",
            "👷 مرتبات الموظفين", "🏠 إيجار المزرعة", "🚛 النقل", "🔧 الصيانة", "📦 أخرى"]
    y = Inches(2.1)
    for c in cats:
        txb(s, c, Inches(9.2), y, Inches(3.7), Inches(0.42),
            font_size=13, color=DARK_TEXT, align=PP_ALIGN.RIGHT, rtl=True)
        y += Inches(0.5)
    reports = [
        ("📊  مخططات شريطية",      "الدخل مقابل المصاريف شهرياً أو سنوياً"),
        ("🔍  تفاصيل تحليلية",      "اضغط على أي شريط لرؤية المعاملات الفردية"),
        ("📋  سجل كامل",            "قوائم المصاريف والدخل مصنفة"),
        ("📁  تصدير CSV",           "مشاركة البيانات كجدول بيانات"),
        ("📄  تصدير PDF",           "تقرير A4 جاهز للطباعة"),
        ("🎯  تتبع الميزانية",       "حدد ميزانية شهرية/سنوية لكل فئة"),
        ("💹  شريط المؤشرات",       "إجمالي الدخل، المصاريف، صافي الربح، أفضل شهر"),
    ]
    y = Inches(1.4)
    for title, desc in reports:
        rect(s, Inches(0.4), y, Inches(8.4), Inches(0.72), GREEN_PALE)
        txb(s, title, Inches(0.5), y + Pt(4), Inches(8.2), Inches(0.38),
            font_size=13, bold=True, color=GREEN_DARK, align=PP_ALIGN.RIGHT, rtl=True)
        txb(s, desc, Inches(0.5), y + Pt(24), Inches(8.2), Inches(0.38),
            font_size=12, color=DARK_TEXT, align=PP_ALIGN.RIGHT, rtl=True)
        y += Inches(0.82)


def slide_buyer_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_MID)
    txb(s, "مميزات المشتري — اكتشف وتسوق المواشي", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=34, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    features = [
        ("🔍  الاكتشاف",         "تصفح المزارع؛ فلتر بالمحافظة، النوع، السعر، التقييم"),
        ("🏡  ملف المزرعة",      "الغلاف، التقييمات، بيانات التواصل، جميع المنتجات"),
        ("📸  إعلانات غنية",     "صور، فيديو، الوزن، العمر، السلالة، الشهادة الصحية"),
        ("❤️  الشهادة الصحية",   "كل حيوان يُظهر حالته الصحية الموثقة"),
        ("📞  تواصل مع البائع",  "واتساب، اتصال مباشر، نسخ الرقم — بنقرة واحدة"),
        ("🛒  تقديم طلب",        "توصيل أو استلام؛ نقداً عند الاستلام أو InstaPay"),
        ("📦  تتبع الطلبات",     "حالة الطلب: معلق ← مؤكد ← في الطريق ← تم التسليم"),
        ("⭐  التقييمات",         "قيّم البائع بعد الاستلام؛ اقرأ تقييمات الآخرين"),
        ("🔔  الإشعارات",        "تنبيهات فورية لتحديثات الطلبات والعروض الخاصة"),
        ("❤️  المفضلة",          "احفظ المزارع المفضلة للوصول السريع"),
    ]
    y = Inches(1.4)
    for i, (feat, desc) in enumerate(features):
        col = 0 if i < 5 else 1
        xi = Inches(0.4) if col == 0 else Inches(6.8)
        row_y = Inches(1.4) + Inches((i % 5) * 1.1)
        rect(s, xi, row_y, Inches(6.1), Inches(0.95), GREEN_PALE)
        txb(s, feat, xi + Inches(0.1), row_y + Pt(4), Inches(5.9), Inches(0.45),
            font_size=13, bold=True, color=GREEN_DARK, align=PP_ALIGN.RIGHT, rtl=True)
        txb(s, desc, xi + Inches(0.1), row_y + Pt(24), Inches(5.9), Inches(0.38),
            font_size=12, color=DARK_TEXT, align=PP_ALIGN.RIGHT, rtl=True)


def slide_coming_soon_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(1.3), GREEN_DARK)
    txb(s, "قريباً", Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
        font_size=34, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    items = [
        ("🔪  الجزار عند الطلب",           "اطلب جزاراً معتمداً يأتي إلى موقعك — بدون تعب"),
        ("🥩  توصيل الذبيحة الجاهزة",      "استلم حيوانك مذبوحاً ومجهزاً أمام بيتك"),
        ("🔨  مزادات المواشي المباشرة",     "تزايد على مواشي مميزة في مزادات أونلاين في الوقت الحقيقي"),
        ("🤖  مساعد تسعير بالذكاء الاصطناعي","توصيات أسعار ذكية بناءً على الوزن والسلالة واتجاهات السوق"),
        ("🌿  إدارة الفروع المتعددة",       "أدر مواقع مزارع متعددة من حساب بائع واحد"),
        ("💳  بوابة الدفع بالبطاقة",        "تكامل كامل مع فوري / PayMob لفيزا وماستركارد"),
    ]
    y = Inches(1.45)
    for icon_title, desc in items:
        rect(s, Inches(0.4), y, Inches(12.5), Inches(0.82), GREEN_PALE)
        txb(s, icon_title, Inches(0.4), y + Pt(5), Inches(12.3), Inches(0.4),
            font_size=15, bold=True, color=GREEN_DARK, align=PP_ALIGN.RIGHT, rtl=True)
        txb(s, desc, Inches(0.4), y + Pt(25), Inches(12.3), Inches(0.38),
            font_size=13, color=DARK_TEXT, align=PP_ALIGN.RIGHT, rtl=True)
        y += Inches(0.92)


def slide_closing_ar(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, GREEN_DARK)
    rect(s, 0, Inches(5.5), SLIDE_W, Inches(2.0), GREEN_MID)
    c = s.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(5), Inches(5))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN_MID; c.line.fill.background()
    txb(s, "شكراً", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
        font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=True)
    txb(s, "فارم فلو — نقل سوق المواشي المصري إلى العصر الرقمي",
        Inches(1), Inches(3.2), Inches(11), Inches(0.8),
        font_size=20, color=RGBColor(0xC8, 0xE6, 0xC9), align=PP_ALIGN.CENTER, italic=True, rtl=True)
    txb(s, "ويب  ·  جوال  ·  سطح المكتب  ·  يعمل أوفلاين",
        Inches(1), Inches(5.8), Inches(11), Inches(0.6),
        font_size=16, color=WHITE, align=PP_ALIGN.CENTER, rtl=True)
    txb(s, "الإصدار 1.0  —  مايو 2026",
        Inches(1), Inches(6.6), Inches(11), Inches(0.5),
        font_size=12, color=RGBColor(0xA5, 0xD6, 0xA7),
        align=PP_ALIGN.CENTER, italic=True, rtl=True)


# ─── Build & Save ────────────────────────────────────────────────

def build_english():
    prs = new_prs()
    slide_cover_en(prs)
    slide_problem_en(prs)
    slide_solution_en(prs)
    slide_platforms_en(prs)
    slide_seller_en(prs)
    slide_herd_en(prs)
    slide_finance_en(prs)
    slide_buyer_en(prs)
    slide_coming_soon_en(prs)
    slide_closing_en(prs)
    out = "FarmFlow_Presentation_EN.pptx"
    prs.save(out)
    print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")


def build_arabic():
    prs = new_prs()
    slide_cover_ar(prs)
    slide_problem_ar(prs)
    slide_solution_ar(prs)
    slide_platforms_ar(prs)
    slide_seller_ar(prs)
    slide_herd_ar(prs)
    slide_finance_ar(prs)
    slide_buyer_ar(prs)
    slide_coming_soon_ar(prs)
    slide_closing_ar(prs)
    out = "FarmFlow_Presentation_AR.pptx"
    prs.save(out)
    print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_english()
    build_arabic()
    print("\nDone! Open both .pptx files in PowerPoint or Google Slides.")
